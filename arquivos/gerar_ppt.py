from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_tcc_presentation():
    # Cria uma nova apresentação
    prs = Presentation()
    
    # Define Layouts Padrões
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Capa ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "IncluiDev"
    subtitle.text = ("Desenvolvimento e Avaliação de uma Plataforma Baseada em Evidências para Apoio ao Ensino de Programação para Pessoas com Deficiência Visual\n\n"
                     "Autores: Lucas R. Pedro & Esteice Janaina Santos Batista\n"
                     "Faculdade de Computação – UFMS\n\nDefesa de TCC - 2026")
                     
    # --- Slide 2: Contexto e Motivação ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Contexto e Motivação"
    tf = body_shape.text_frame
    tf.text = "O ensino de programação é fortemente dependente de metáforas visuais (diagramas, blocos, indentação)."
    p = tf.add_paragraph()
    p.text = "Isso gera intensa sobrecarga cognitiva para estudantes cegos utilizando leitores de tela."
    p = tf.add_paragraph()
    p.text = "Existem excelentes tecnologias assistivas acadêmicas, mas elas estão dispersas pela web."
    p = tf.add_paragraph()
    p.text = "A fragmentação dificulta o trabalho de educadores e alunos no planejamento de ensino."

    # --- Slide 3: O Problema e Objetivos ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Problema de Pesquisa e Objetivos"
    tf = body_shape.text_frame
    tf.text = "Questão: Como centralizar, classificar e recomendar de forma acessível essas soluções educacionais?"
    p = tf.add_paragraph()
    p.text = "Objetivo Geral: Construir o repositório 'IncluiDev'."
    p = tf.add_paragraph()
    p.text = "Objetivos Específicos:"
    p = tf.add_paragraph()
    p.text = "Organizar recursos em uma taxonomia"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Garantir acessibilidade (WCAG 2.1 AA)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Criar o 'Assistente Guiado' (Busca Orientada)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Validar a usabilidade e a arquitetura"
    p.level = 1

    # --- Slide 4: Fundamentação Teórica ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Fundamentação Teórica"
    tf = body_shape.text_frame
    tf.text = "As respostas da comunidade para o ensino acessível convergem em:"
    p = tf.add_paragraph()
    p.text = "Abordagens Tangíveis (Hardware)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Retirar o aluno da tela. Ex: Code Jumper."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Abordagens de Software"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Linguagens baseadas em evidências de percepção sonora (Ex: Quorum)."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Essas duas grandes abordagens foram catalogadas como recursos no repositório."

    # --- Slide 5: A Solução ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "A Solução: Portal IncluiDev"
    tf = body_shape.text_frame
    tf.text = "Arquitetura: Single Page Application em React e Node.js."
    p = tf.add_paragraph()
    p.text = "Acervo inicial de 34 recursos fundamentais (Papers, ferramentas, normas técnicas)."
    p = tf.add_paragraph()
    p.text = "Desenvolvimento com 'Acessibilidade Nativa'."
    p = tf.add_paragraph()
    p.text = "Uso de atributos dinâmicos WAI-ARIA visando compatibilidade máxima com NVDA."
    p.level = 1

    # --- Slide 6: O Assistente Guiado ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "A Inovação: O Assistente Guiado"
    tf = body_shape.text_frame
    tf.text = "O usuário de leitor de tela não precisa ler longas tabelas de dados."
    p = tf.add_paragraph()
    p.text = "Busca Dinâmica por Perfil de Usuário:"
    p = tf.add_paragraph()
    p.text = "O sistema pergunta se o usuário é Estudante, Professor, Pesquisador ou Desenvolvedor."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Com mais 2 perguntas sobre o objetivo, o Motor de Inferência cruza tags semânticas."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Resultado limpo, contendo apenas recomendações perfeitamente aderentes à necessidade."
    p.level = 1

    # --- Slide 7: Validação de Acessibilidade ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Validação 1: Acessibilidade Técnica"
    tf = body_shape.text_frame
    tf.text = "Auditoria Automatizada:"
    p = tf.add_paragraph()
    p.text = "Lighthouse: Score 96/100 para Acessibilidade Desktop/Mobile."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "WAVE: Nenhum erro estrutural crítico."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Auditoria Manual:"
    p = tf.add_paragraph()
    p.text = "Navegação fluida sem 'Keyboard Traps'."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uso de aria-live='assertive' notifica as mutações da tela instantaneamente para o leitor de tela."
    p.level = 1

    # --- Slide 8: Validação por Simulação ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Validação 2: Simulação por Agentes de IA"
    tf = body_shape.text_frame
    tf.text = "Avaliação Heurística Preliminar (Metodologia embasada em Seshadri et al., 2026)."
    p = tf.add_paragraph()
    p.text = "Utilização de Grandes Modelos de Linguagem (LLMs) interpretando as Personas."
    p = tf.add_paragraph()
    p.text = "O LLM foi exposto à representação textual das telas e forçado a preencher a System Usability Scale (SUS) e o TAM."
    p = tf.add_paragraph()
    p.text = "A IA atuou como um 'Filtro de Qualidade' preparatório para garantir que a lógica das tags faz sentido."

    # --- Slide 9: Resultados da Simulação ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Resultados da Validação Lógica"
    tf = body_shape.text_frame
    tf.text = "Todos os perfis completaram a jornada guiada sem bloqueios ou rotas vazias."
    p = tf.add_paragraph()
    p.text = "A utilidade percebida (TAM) na geração das dissertações foi altíssima:"
    p = tf.add_paragraph()
    p.text = "Estudantes reportaram não precisarem clicar por dezenas de links irrelevantes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Pesquisadores consideraram a interface mais ágil que buscas booleanas."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Nota Metodológica: Assumimos os dados quantitativos apenas como validação de fluxo lógico prévio, cientes de que IAs não sentem a fadiga sensorial orgânica."

    # --- Slide 10: Conclusão ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusão"
    tf = body_shape.text_frame
    tf.text = "O IncluiDev consolida, de fato, a teoria fragmentada num portal acessível."
    p = tf.add_paragraph()
    p.text = "O Assistente Guiado atingiu seu propósito heurístico de reduzir a carga cognitiva, fornecendo recomendação baseada em necessidade."
    p = tf.add_paragraph()
    p.text = "O desenvolvimento do sistema manteve rígida conformidade aos padrões estabelecidos pela W3C (WCAG 2.1)."

    # --- Slide 11: Limitações e Trabalhos Futuros ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Limitações e Trabalhos Futuros"
    tf = body_shape.text_frame
    tf.text = "Limitação Principal: Os testes foram baseados na simulação de agentes de IA."
    p = tf.add_paragraph()
    p.text = "Trabalhos Futuros:"
    p = tf.add_paragraph()
    p.text = "Execução de testes empíricos com grupos reais (Estudantes PcDV e Educadores)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Expansão da base de recursos através de curadoria colaborativa."
    p.level = 1
    
    # --- Slide 12: Agradecimentos ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Obrigado!"
    subtitle.text = "Perguntas?"

    # Salva a apresentação
    prs.save("docs/Apresentacao_TCC_IncluiDev.pptx")
    print("Apresentação salva com sucesso em docs/Apresentacao_TCC_IncluiDev.pptx")

if __name__ == '__main__':
    create_tcc_presentation()